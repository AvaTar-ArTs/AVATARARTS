"""
PPTX Builder - create editable PPTX files.
"""
import os
import logging
import base64
import io
from typing import List, Dict, Any, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# Try to import python-pptx, fall back if unavailable.
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RgbColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx is not installed; PPTX export is unavailable")

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class PPTXBuilder:
    """Build editable PPTX files from structured content."""

    # Standard slide size (16:9).
    DEFAULT_SLIDE_WIDTH_INCHES = 13.333
    DEFAULT_SLIDE_HEIGHT_INCHES = 7.5

    # Default DPI.
    DEFAULT_DPI = 96

    def __init__(self, slide_width_inches: float = None, slide_height_inches: float = None):
        """
        Initialize the PPTX builder.

        Args:
            slide_width_inches: Slide width in inches.
            slide_height_inches: Slide height in inches.
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Run pip install python-pptx.")

        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs = None
        self.current_slide = None

    def create_presentation(self) -> 'Presentation':
        """Create a new presentation."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs

    def add_blank_slide(self):
        """Add a blank slide."""
        if not self.prs:
            self.create_presentation()

        # Use blank layout (usually layout 6).
        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide

    def add_title_slide(self, title: str, subtitle: str = None):
        """
        Add a title slide.

        Args:
            title: Main title.
            subtitle: Subtitle.
        """
        if not self.prs:
            self.create_presentation()

        # Use title layout.
        title_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_layout)

        # Set title.
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set subtitle.
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        self.current_slide = slide
        return slide

    def add_content_slide(self, title: str, content: str, bullet_points: List[str] = None):
        """
        Add a content slide.

        Args:
            title: Slide title.
            content: Body text.
            bullet_points: Bullet points.
        """
        if not self.prs:
            self.create_presentation()

        # Use title-and-content layout.
        content_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(content_layout)

        # Set title.
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set content.
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame

            if bullet_points:
                for i, point in enumerate(bullet_points):
                    if i == 0:
                        tf.text = point
                    else:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
            elif content:
                tf.text = content

        self.current_slide = slide
        return slide

    def add_image_to_slide(self, slide, image_data: str,
                           left: float = 0, top: float = 0,
                           width: float = None, height: float = None):
        """
        Add an image to a slide.

        Args:
            slide: Target slide.
            image_data: Base64-encoded image data.
            left: Left margin in inches.
            top: Top margin in inches.
            width: Width in inches.
            height: Height in inches.
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL is not installed; cannot process images")
            return

        try:
            # Decode Base64 image.
            if image_data.startswith('data:image'):
                # Remove data URL prefix.
                image_data = image_data.split(',')[1]

            image_bytes = base64.b64decode(image_data)
            image_stream = io.BytesIO(image_bytes)

            # Use full slide size if not specified.
            if width is None:
                width = self.slide_width_inches
            if height is None:
                height = self.slide_height_inches

            # Add image.
            slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height)
            )
            logger.debug("Image added to slide")

        except Exception as e:
            logger.error(f"Failed to add image: {e}")

    def add_text_box(self, slide, text: str,
                     left: float, top: float,
                     width: float, height: float,
                     font_size: int = 18,
                     bold: bool = False,
                     align: str = 'left'):
        """
        Add a text box.

        Args:
            slide: Target slide.
            text: Text content.
            left, top, width, height: Position and size (inches).
            font_size: Font size (points).
            bold: Whether to use bold text.
            align: Alignment ('left', 'center', 'right').
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.text = text
        tf.word_wrap = True

        # Set paragraph formatting.
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold

            if align == 'center':
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT

        return textbox

    def build_from_presentation_data(self, presentation_data: Dict[str, Any]) -> 'Presentation':
        """
        Build a PPTX from presentation data.

        Args:
            presentation_data: Presentation data dict, including title and slides.

        Returns:
            Presentation instance.
        """
        self.create_presentation()

        slides = presentation_data.get('slides', [])

        for slide_data in slides:
            slide_type = slide_data.get('type', 'content')
            title = slide_data.get('title', '')
            content = slide_data.get('content', '')
            image = slide_data.get('image')

            if slide_type == 'title':
                # Title slide.
                slide = self.add_blank_slide()

                # Add background image if present.
                if image:
                    self.add_image_to_slide(slide, image)

                # Add title text.
                self.add_text_box(
                    slide, title,
                    left=0.5, top=2.5,
                    width=12.333, height=1.5,
                    font_size=44, bold=True, align='center'
                )

                # Add subtitle.
                if content:
                    self.add_text_box(
                        slide, content,
                        left=0.5, top=4.2,
                        width=12.333, height=1,
                        font_size=24, align='center'
                    )
            else:
                # Content slide.
                slide = self.add_blank_slide()

                # Add background image if present.
                if image:
                    self.add_image_to_slide(slide, image)

                # Add title.
                self.add_text_box(
                    slide, title,
                    left=0.5, top=0.3,
                    width=12.333, height=0.8,
                    font_size=32, bold=True
                )

                # Add content.
                if content:
                    # Parse bullet content.
                    lines = content.split('\n')
                    formatted_content = '\n'.join(lines)

                    self.add_text_box(
                        slide, formatted_content,
                        left=0.5, top=1.3,
                        width=12.333, height=5.5,
                        font_size=20
                    )

        return self.prs

    def save(self, output_path: str):
        """
        Save the presentation to a file.

        Args:
            output_path: Output file path.
        """
        if not self.prs:
            raise ValueError("No presentation to save")

        # Ensure the directory exists.
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)

        self.prs.save(output_path)
        logger.info(f"Presentation saved to: {output_path}")

    def save_to_bytes(self) -> bytes:
        """
        Save the presentation to bytes.

        Returns:
            PPTX bytes.
        """
        if not self.prs:
            raise ValueError("No presentation to save")

        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def save_to_base64(self) -> str:
        """
        Save the presentation as Base64.

        Returns:
            Base64-encoded PPTX data.
        """
        pptx_bytes = self.save_to_bytes()
        return base64.b64encode(pptx_bytes).decode('utf-8')


def is_pptx_available() -> bool:
    """Check whether PPTX support is available."""
    return PPTX_AVAILABLE
